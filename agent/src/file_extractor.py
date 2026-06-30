"""
Multi-format text extractor for DLP file scanning.
Returns plain text for the classifier or None when the file should be skipped.
"""

from pathlib import Path
from loguru import logger

# Per-document extraction cap — generous enough to cover full files but prevents
# sending massive strings to the classifier.
_EXTRACT_LIMIT = 50_000  # characters

_PLAIN_TEXT_EXTS = frozenset({
    ".txt", ".csv", ".json", ".log", ".xml", ".md",
    ".yaml", ".yml", ".ini", ".cfg", ".toml",
    ".py", ".js", ".ts", ".jsx", ".tsx",
    ".html", ".htm", ".css", ".sql",
    ".sh", ".bat", ".ps1",
})


def extract(file_path: str) -> str | None:
    """
    Extract readable text from *file_path* based on its extension.

    Returns:
        str  — extracted text (may be truncated to _EXTRACT_LIMIT chars)
        None — file should be skipped (unsupported, failed, or empty)
    """
    path = Path(file_path)
    ext  = path.suffix.lower()

    try:
        if ext in _PLAIN_TEXT_EXTS:
            return _plain_text(path)
        if ext == ".docx":
            return _docx(path)
        if ext == ".xlsx":
            return _xlsx(path)
        if ext == ".pptx":
            return _pptx(path)
        if ext == ".pdf":
            return _pdf(path)
        # Unknown extension — attempt UTF-8 text read, silently give up if binary
        return _plain_text(path, silent=True)
    except Exception as exc:
        logger.warning(f"[EXTRACTOR] Extraction failed for '{path.name}': {exc}")
        return None


# ── Plain text ────────────────────────────────────────────────────────────────

def _plain_text(path: Path, silent: bool = False) -> str | None:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
        return text[:_EXTRACT_LIMIT]
    except (PermissionError, OSError) as exc:
        if not silent:
            logger.warning(f"[EXTRACTOR] Cannot read '{path.name}': {exc}")
        return None
    except Exception as exc:
        if not silent:
            logger.warning(f"[EXTRACTOR] Text read error '{path.name}': {exc}")
        return None


# ── Microsoft Word (.docx) ────────────────────────────────────────────────────

def _docx(path: Path) -> str | None:
    from docx import Document  # python-docx
    doc   = Document(str(path))
    parts: list[str] = []

    for para in doc.paragraphs:
        t = para.text.strip()
        if t:
            parts.append(t)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                t = cell.text.strip()
                if t:
                    parts.append(t)

    text = "\n".join(parts)
    logger.debug(f"[EXTRACTOR] .docx '{path.name}': {len(parts)} blocks, {len(text)} chars")
    return text[:_EXTRACT_LIMIT]


# ── Microsoft Excel (.xlsx) ───────────────────────────────────────────────────

def _xlsx(path: Path) -> str | None:
    import openpyxl
    wb    = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = []

    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_vals = [str(v) for v in row if v is not None]
            if row_vals:
                parts.append("  ".join(row_vals))

    n_sheets = len(wb.worksheets)
    wb.close()
    text = "\n".join(parts)
    logger.debug(f"[EXTRACTOR] .xlsx '{path.name}': {n_sheets} sheet(s), {len(text)} chars")
    return text[:_EXTRACT_LIMIT]


# ── Microsoft PowerPoint (.pptx) ─────────────────────────────────────────────

def _pptx(path: Path) -> str | None:
    from pptx import Presentation  # python-pptx
    prs   = Presentation(str(path))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())

    text = "\n".join(parts)
    logger.debug(f"[EXTRACTOR] .pptx '{path.name}': {len(prs.slides)} slide(s), {len(text)} chars")
    return text[:_EXTRACT_LIMIT]


# ── PDF ───────────────────────────────────────────────────────────────────────

def _pdf(path: Path) -> str | None:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    parts: list[str] = []

    for page in reader.pages:
        try:
            t = page.extract_text()
            if t and t.strip():
                parts.append(t)
        except Exception:
            pass  # skip unreadable pages

    text = "\n".join(parts)
    logger.debug(f"[EXTRACTOR] .pdf '{path.name}': {len(reader.pages)} page(s), {len(text)} chars")
    return text[:_EXTRACT_LIMIT]
