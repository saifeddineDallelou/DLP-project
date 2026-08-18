from unittest.mock import patch, MagicMock

import file_extractor


def test_extract_plain_text(tmp_path):
    f = tmp_path / "notes.txt"
    f.write_text("SSN: 123-45-6789", encoding="utf-8")
    assert file_extractor.extract(str(f)) == "SSN: 123-45-6789"


def test_extract_truncates_to_limit(tmp_path):
    f = tmp_path / "big.txt"
    f.write_text("a" * 60_000, encoding="utf-8")
    result = file_extractor.extract(str(f))
    assert len(result) == file_extractor._EXTRACT_LIMIT


def test_extract_missing_file_returns_none(tmp_path):
    missing = tmp_path / "nope.txt"
    assert file_extractor.extract(str(missing)) is None


def test_extract_unknown_extension_falls_back_to_text(tmp_path):
    f = tmp_path / "weird.xyz"
    f.write_text("plain content here", encoding="utf-8")
    assert file_extractor.extract(str(f)) == "plain content here"


def test_extract_unknown_extension_binary_is_silently_skipped(tmp_path):
    f = tmp_path / "blob.xyz"
    f.write_bytes(b"\xff\xfe\x00\x01binary junk")
    # Should not raise; utf-8 with errors="replace" always returns *some* text,
    # so this exercises the silent-fallback path without crashing.
    result = file_extractor.extract(str(f))
    assert result is None or isinstance(result, str)


def test_extract_docx(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("Card number: 4111111111111111")
    table = doc.add_table(rows=1, cols=1)
    table.rows[0].cells[0].text = "table cell text"
    path = tmp_path / "doc.docx"
    doc.save(str(path))

    text = file_extractor.extract(str(path))
    assert "Card number: 4111111111111111" in text
    assert "table cell text" in text


def test_extract_xlsx(tmp_path):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "email"
    ws["B1"] = "a@b.com"
    path = tmp_path / "book.xlsx"
    wb.save(str(path))

    text = file_extractor.extract(str(path))
    assert "[Sheet: Sheet1]" in text
    assert "email" in text
    assert "a@b.com" in text


def test_extract_pptx(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "confidential budget numbers"
    path = tmp_path / "deck.pptx"
    prs.save(str(path))

    text = file_extractor.extract(str(path))
    assert "[Slide 1]" in text
    assert "confidential budget numbers" in text


def test_extract_pdf_uses_pypdf(tmp_path):
    path = tmp_path / "doc.pdf"
    path.write_bytes(b"%PDF-1.4 fake")  # content irrelevant, PdfReader is mocked

    fake_page = MagicMock()
    fake_page.extract_text.return_value = "extracted pdf text"
    fake_reader = MagicMock()
    fake_reader.pages = [fake_page]

    with patch("pypdf.PdfReader", return_value=fake_reader):
        text = file_extractor.extract(str(path))

    assert text == "extracted pdf text"


def test_extract_pdf_skips_unreadable_pages(tmp_path):
    path = tmp_path / "doc.pdf"
    path.write_bytes(b"%PDF-1.4 fake")

    bad_page = MagicMock()
    bad_page.extract_text.side_effect = Exception("corrupt page")
    good_page = MagicMock()
    good_page.extract_text.return_value = "readable text"
    fake_reader = MagicMock()
    fake_reader.pages = [bad_page, good_page]

    with patch("pypdf.PdfReader", return_value=fake_reader):
        text = file_extractor.extract(str(path))

    assert text == "readable text"
